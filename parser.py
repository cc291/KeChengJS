from __future__ import annotations

from io import BytesIO
from typing import BinaryIO

from docx import Document
from pptx import Presentation


def _to_binary_stream(file: bytes | BinaryIO) -> BytesIO | BinaryIO:
    if isinstance(file, bytes):
        return BytesIO(file)
    return file


def parse_ppt(file: bytes | BinaryIO) -> str:
    """
    提取 PPT 全部幻灯片文本，并按页面顺序拼接。
    过滤少于 5 个字符的碎片文本。
    """
    stream = _to_binary_stream(file)
    presentation = Presentation(stream)

    slide_texts: list[str] = []
    for slide in presentation.slides:
        fragments: list[str] = []
        for shape in slide.shapes:
            if not getattr(shape, "has_text_frame", False):
                continue
            raw_text = shape.text.strip()
            if len(raw_text) < 5:
                continue
            fragments.append(raw_text)

        if fragments:
            slide_texts.append("\n".join(fragments))

    return "\n\n".join(slide_texts).strip()


def parse_docx(file: bytes | BinaryIO) -> str:
    """
    提取 Word 教案全文，包含段落和表格内容。
    针对工科教案优化：递归提取表格单元格中的文本。
    """
    stream = _to_binary_stream(file)
    doc = Document(stream)
    full_content = []

    # 1. 提取所有段落文本
    for para in doc.paragraphs:
        text = para.text.strip()
        if text:
            full_content.append(text)

    # 2. 核心优化：提取所有表格中的文本 [cite: 135, 141]
    # 工科教案的知识点和思政映射通常写在表格单元格中
    for table in doc.tables:
        for row in table.rows:
            row_text = []
            for cell in row.cells:
                # 提取单元格内所有段落并合并
                cell_text = " ".join([p.text.strip() for p in cell.paragraphs if p.text.strip()])
                if cell_text:
                    row_text.append(cell_text)
            if row_text:
                # 将一行内的内容用 | 分隔，模拟表格结构利于 AI 理解
                full_content.append(" | ".join(row_text))

    return "\n".join(full_content).strip()